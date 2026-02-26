#!/usr/bin/env python
"""Build Skills Cookbook tutorial deck from an Imperial UK DRI POTX.

Generates a short intro section covering privacy and model/provider choices,
using the provided PowerPoint template for branding.
"""

from __future__ import annotations

import argparse
from pathlib import Path
import zipfile
from xml.etree import ElementTree as ET

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.oxml.ns import qn


DEFAULT_TEMPLATE = Path(
    "/Users/jaymoore/Documents/JAY_PhD/imperial/Support/IMPERIAL UK DRI - PPT Template - 2025.potx"
)


_PRESENTATION_CT = (
    "application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml"
)
_TEMPLATE_CT = (
    "application/vnd.openxmlformats-officedocument.presentationml.template.main+xml"
)


def _materialize_potx_as_pptx(potx_path: Path, pptx_path: Path) -> Path:
    """Convert a .potx to a .pptx by adjusting the content type.

    python-pptx refuses to open template packages (.potx) directly.
    The underlying package structure is compatible; we just need the
    presentation.xml content-type to be the .pptx variant.
    """

    pptx_path.parent.mkdir(parents=True, exist_ok=True)

    with (
        zipfile.ZipFile(potx_path, "r") as zin,
        zipfile.ZipFile(pptx_path, "w") as zout,
    ):
        for info in zin.infolist():
            data = zin.read(info.filename)

            if info.filename == "[Content_Types].xml":
                root = ET.fromstring(data)
                # The namespace is usually the default; ElementTree preserves it.
                for override in root.findall("{*}Override"):
                    part_name = override.get("PartName")
                    if (
                        part_name == "/ppt/presentation.xml"
                        and override.get("ContentType") == _TEMPLATE_CT
                    ):
                        override.set("ContentType", _PRESENTATION_CT)
                data = ET.tostring(root, encoding="utf-8", xml_declaration=True)

            zout.writestr(info, data)

    return pptx_path


def _add_title(
    slide,
    title: str,
    subtitle: str | None = None,
    *,
    x=Inches(0.75),
    y=Inches(0.8),
    w=Inches(7.55),
    h=Inches(1.4),
):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title
    run.font.name = "Calibri Light"
    run.font.size = Pt(46)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x00, 0x00, 0x00)

    if subtitle:
        p2 = tf.add_paragraph()
        p2.space_before = Pt(8)
        run2 = p2.add_run()
        run2.text = subtitle
        run2.font.name = "Calibri"
        run2.font.size = Pt(22)
        run2.font.color.rgb = RGBColor(0x33, 0x33, 0x33)


def _add_bullets(
    slide,
    title: str,
    bullets: list[str],
    *,
    x=Inches(0.75),
    y=Inches(0.9),
    w=Inches(7.55),
    h=Inches(6.2),
):
    _add_title(slide, title, None, x=x, y=y, w=w, h=Inches(0.9))

    box = slide.shapes.add_textbox(x, y + Inches(1.0), w, h - Inches(1.0))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True

    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b
        p.level = 0
        p.font.name = "Calibri"
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(0x11, 0x11, 0x11)
        p.space_after = Pt(6)


def _add_bullet_with_link(
    tf,
    text_before: str,
    link_text: str,
    url: str,
    text_after: str = "",
    *,
    font_size=Pt(22),
):
    p = tf.add_paragraph()
    p.level = 0
    p.space_after = Pt(6)

    r1 = p.add_run()
    r1.text = text_before
    r1.font.name = "Calibri"
    r1.font.size = font_size
    r1.font.color.rgb = RGBColor(0x11, 0x11, 0x11)

    r2 = p.add_run()
    r2.text = link_text
    r2.hyperlink.address = url
    r2.font.name = "Calibri"
    r2.font.size = font_size
    r2.font.underline = True
    r2.font.color.rgb = RGBColor(0x00, 0x4C, 0x8A)

    if text_after:
        r3 = p.add_run()
        r3.text = text_after
        r3.font.name = "Calibri"
        r3.font.size = font_size
        r3.font.color.rgb = RGBColor(0x11, 0x11, 0x11)


def build_deck(template_path: Path, out_path: Path) -> None:
    if not template_path.exists():
        raise FileNotFoundError(f"Template not found: {template_path}")

    template_for_pptx = template_path
    if template_path.suffix.lower() == ".potx":
        template_for_pptx = _materialize_potx_as_pptx(
            template_path,
            Path("slides/.build") / (template_path.stem + ".pptx"),
        )

    prs = Presentation(str(template_for_pptx))

    # The provided template includes example/instructional slides.
    # Remove them so the output deck contains only our tutorial slides.
    sldIdLst = prs.slides._sldIdLst  # pylint: disable=protected-access
    for sldId in list(sldIdLst):
        rId = sldId.get(qn("r:id"))
        prs.part.drop_rel(rId)
        sldIdLst.remove(sldId)

    # The template has a right-side logo area; keep content within this safe zone.
    x0 = Inches(0.75)
    safe_w = Inches(7.55)  # right edge ~ 8.30in, avoids template logos

    layout = prs.slide_layouts[0]

    # Slide 1: Title
    s1 = prs.slides.add_slide(layout)
    _add_title(
        s1,
        "Skills Cookbook",
        "Agentic work and coding assistants for the lab",
        x=x0,
        y=Inches(0.9),
        w=safe_w,
        h=Inches(1.6),
    )

    # Slide 2: Privacy
    s2 = prs.slides.add_slide(layout)
    _add_title(
        s2, "Privacy & Data Handling", x=x0, y=Inches(0.8), w=safe_w, h=Inches(0.9)
    )
    box = s2.shapes.add_textbox(x0, Inches(1.8), safe_w, Inches(5.8))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True

    p0 = tf.paragraphs[0]
    p0.text = "Treat prompts like sharing with a third party"
    p0.level = 0
    p0.font.name = "Calibri"
    p0.font.size = Pt(26)
    p0.font.color.rgb = RGBColor(0x11, 0x11, 0x11)
    p0.space_after = Pt(8)

    for b in [
        "Avoid secrets (API keys, tokens, credentials)",
        "Prefer de-identified snippets; paste the minimum needed",
        "Use provider privacy settings and retention policies",
    ]:
        p = tf.add_paragraph()
        p.text = b
        p.level = 0
        p.font.name = "Calibri"
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(0x11, 0x11, 0x11)
        p.space_after = Pt(6)

    _add_bullet_with_link(
        tf,
        "OpenRouter privacy settings: ",
        "openrouter.ai/settings/privacy",
        "https://openrouter.ai/settings/privacy",
    )
    _add_bullet_with_link(
        tf,
        "OpenRouter retention (ZDR): ",
        "openrouter.ai/docs/guides/features/zdr",
        "https://openrouter.ai/docs/guides/features/zdr#openrouters-retention-policy",
    )

    # Slide 3: Providers & Models
    s3 = prs.slides.add_slide(layout)
    _add_bullets(
        s3,
        "Providers & Models",
        [
            "OpenCode = the interface (agent workflow on your machine)",
            "OpenRouter = provider (one account to access many models)",
            "Course default: Kimi K2.5 via OpenRouter",
            "Opus 4.5 option: OpenCode → provider GitHub Copilot (GitHub Education)",
        ],
        x=x0,
        y=Inches(0.8),
        w=safe_w,
        h=Inches(6.6),
    )

    # Slide 4: Context & Cost
    s4 = prs.slides.add_slide(layout)
    _add_title(
        s4, "Models, Context, and Cost", x=x0, y=Inches(0.8), w=safe_w, h=Inches(0.9)
    )
    box4 = s4.shapes.add_textbox(x0, Inches(1.7), safe_w, Inches(1.2))
    tf4 = box4.text_frame
    tf4.clear()
    tf4.word_wrap = True
    for i, b in enumerate(
        [
            "Context window = conversation + files + tool output, resent each message",
            "Keep context small: /clear between tasks, /compact mid-task",
        ]
    ):
        p = tf4.paragraphs[0] if i == 0 else tf4.add_paragraph()
        p.text = b
        p.level = 0
        p.font.name = "Calibri"
        p.font.size = Pt(22)
        p.font.color.rgb = RGBColor(0x11, 0x11, 0x11)
        p.space_after = Pt(4)

    img = Path("assets/images/context-cost-per-message.png")
    if img.exists():
        s4.shapes.add_picture(str(img), x0, Inches(2.8), width=safe_w)

    # Slide 5: Agenda (quick orientation)
    s5 = prs.slides.add_slide(layout)
    _add_bullets(
        s5,
        "Today",
        [
            "Setup: OpenCode + provider + model selection",
            "How to prompt in Plan vs Build mode",
            "Three short demos: Penguins → Slides → Mini app",
            "Context management + file-based workflows (PLAN.md / GSD)",
        ],
        x=x0,
        y=Inches(0.8),
        w=safe_w,
        h=Inches(6.6),
    )

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))


def main() -> int:
    ap = argparse.ArgumentParser()
    ap.add_argument("--template", type=Path, default=DEFAULT_TEMPLATE)
    ap.add_argument(
        "--out",
        type=Path,
        default=Path("slides/skills-cookbook-tutorial.pptx"),
    )
    args = ap.parse_args()
    build_deck(args.template, args.out)
    print(f"Wrote: {args.out}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
